import os
import logging
import mimetypes
from pathlib import Path
from typing import Optional, Dict, Any

# Data processing libraries
import pandas as pd
from bs4 import BeautifulSoup
from pypdf import PdfReader
import docx
from pptx import Presentation
import ebooklib
from ebooklib import epub

# Image & OCR
from PIL import Image
import pytesseract

# Utils
import chardet

logger = logging.getLogger("ark_monster")

class MonsterLoader:
    """
    Універсальний завантажувач контенту.
    Підтримує: TXT, MD, CSV, JSON, XML, HTML, PDF, DOCX, XLSX, PPTX, EPUB, 
    та зображення (JPG, PNG, TIFF, WEBP, BMP) через OCR.
    """

    def __init__(self, ocr_enabled: bool = True, ocr_lang: str = 'ukr+eng'):
        self.ocr_enabled = ocr_enabled
        self.ocr_lang = ocr_lang
        
        # Mapping extensions to handlers
        self._handlers = {
            # Text / Code
            '.txt': self._read_text,
            '.md': self._read_text,
            '.json': self._read_text,
            '.xml': self._read_text,
            '.csv': self._read_csv,
            '.py': self._read_text,
            '.js': self._read_text,
            
            # Web / Books
            '.html': self._read_html,
            '.htm': self._read_html,
            '.epub': self._read_epub,
            
            # PDF
            '.pdf': self._read_pdf,
            
            # Office
            '.docx': self._read_docx,
            '.doc': self._read_binary_doc_warning, # Old doc format is hard
            '.xlsx': self._read_excel,
            '.xls': self._read_excel,
            '.pptx': self._read_pptx,
            
            # Images
            '.jpg': self._read_image,
            '.jpeg': self._read_image,
            '.png': self._read_image,
            '.bmp': self._read_image,
            '.tiff': self._read_image,
            '.webp': self._read_image,
        }

    def load(self, file_path: Path) -> Optional[str]:
        """Головний метод маршрутизації."""
        if not file_path.exists():
            logger.error(f"File not found: {file_path}")
            return None

        suffix = file_path.suffix.lower()
        handler = self._handlers.get(suffix)

        try:
            if handler:
                logger.info(f"Processing {file_path.name} as {suffix}")
                return handler(file_path)
            else:
                # Fallback: спробувати прочитати як текст
                logger.warning(f"Unknown extension {suffix}, trying text fallback: {file_path}")
                return self._read_text(file_path)
        except Exception as e:
            logger.error(f"Failed to process {file_path}: {e}")
            return None

    # --- HANDLERS ---

    def _read_text(self, path: Path) -> str:
        """Читає текстові файли з автовизначенням кодування."""
        raw_data = path.read_bytes()
        detection = chardet.detect(raw_data)
        encoding = detection['encoding'] or 'utf-8'
        
        try:
            return raw_data.decode(encoding)
        except UnicodeDecodeError:
            return raw_data.decode('utf-8', errors='ignore')

    def _read_csv(self, path: Path) -> str:
        """Перетворює CSV в Markdown таблицю."""
        try:
            df = pd.read_csv(path)
            # Якщо таблиця велика, беремо JSON, якщо мала - Markdown
            if len(df) > 50:
                return df.to_json(orient='records', force_ascii=False)
            return df.to_markdown(index=False)
        except Exception:
            return self._read_text(path)

    def _read_html(self, path: Path) -> str:
        with open(path, 'rb') as f:
            soup = BeautifulSoup(f, 'html.parser')
            # Видаляємо скрипти та стилі
            for script in soup(["script", "style"]):
                script.extract()
            return soup.get_text(separator='\n')

    def _read_pdf(self, path: Path) -> str:
        text = []
        try:
            reader = PdfReader(path)
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
                else:
                    # TODO: Якщо текст пустий, тут можна викликати OCR для сторінки (image extraction)
                    text.append(f"[Page {i+1}: No text layer detected]")
            return "\n\n".join(text)
        except Exception as e:
            logger.error(f"PDF Error: {e}")
            return ""

    def _read_docx(self, path: Path) -> str:
        doc = docx.Document(path)
        full_text = []
        
        # Параграфи
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        
        # Таблиці
        for table in doc.tables:
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                full_text.append(" | ".join(row_data))
        
        return "\n".join(full_text)

    def _read_pptx(self, path: Path) -> str:
        prs = Presentation(path)
        text_runs = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)

    def _read_excel(self, path: Path) -> str:
        # Читаємо всі листи
        dfs = pd.read_excel(path, sheet_name=None)
        output = []
        for sheet_name, df in dfs.items():
            output.append(f"### Sheet: {sheet_name}")
            if len(df) > 50:
                output.append(df.to_json(orient='records', force_ascii=False))
            else:
                output.append(df.to_markdown(index=False))
        return "\n\n".join(output)

    def _read_epub(self, path: Path) -> str:
        book = epub.read_epub(path)
        chapters = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                chapters.append(soup.get_text(separator='\n'))
        return "\n\n".join(chapters)

    def _read_image(self, path: Path) -> str:
        """
        Використовує OCR (Tesseract) для витягування тексту з картинки.
        """
        if not self.ocr_enabled:
            return f"[Image: {path.name}] (OCR disabled)"

        try:
            image = Image.open(path)
            # Конвертуємо інфо про картинку в текст
            meta_info = f"[Image Metadata: {image.format}, {image.size}, {image.mode}]"
            
            # Виконуємо OCR
            extracted_text = pytesseract.image_to_string(image, lang=self.ocr_lang)
            
            if not extracted_text.strip():
                return f"{meta_info}\n[No text detected via OCR]"
            
            return f"{meta_info}\n=== OCR RESULT ===\n{extracted_text}"
            
        except Exception as e:
            logger.error(f"OCR Error on {path}: {e}")
            return f"[Image processing failed: {e}]"

    def _read_binary_doc_warning(self, path: Path) -> str:
        return "[WARNING: Old .doc format not supported directly. Convert to .docx]"
